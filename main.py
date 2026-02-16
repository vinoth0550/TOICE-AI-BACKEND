

# main.py

from fastapi import FastAPI, UploadFile, File, Form
from fastapi.responses import JSONResponse
import os
import shutil
from datetime import datetime
from gemini_service import generate_prd, transcribe_audio
from database import save_prd

from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation

app = FastAPI()

UPLOAD_DIR = "uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)


# ---------------------------
# FILE TEXT EXTRACTION
# ---------------------------

def extract_text_from_docx(path):
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs])


def extract_text_from_pdf(path):
    reader = PdfReader(path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text


def extract_text_from_ppt(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text



# MAIN PRD ENDPOINT


@app.post("/generate-prd")
async def generate_prd_endpoint(
    text: str = Form(None),
    file: UploadFile = File(None)
):

    raw_input_text = ""
    transcript = None
    input_type = None

    # 1️⃣ If Text Provided
    if text:
        raw_input_text = text
        input_type = "text"

    # 2️⃣ If File Provided
    if file:
        file_path = os.path.join(UPLOAD_DIR, file.filename)

        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        if file.filename.endswith(".docx"):
            raw_input_text = extract_text_from_docx(file_path)
            input_type = "docx"

        elif file.filename.endswith(".pdf"):
            raw_input_text = extract_text_from_pdf(file_path)
            input_type = "pdf"

        elif file.filename.endswith(".pptx"):
            raw_input_text = extract_text_from_ppt(file_path)
            input_type = "ppt"

        elif file.filename.endswith(".wav") or file.filename.endswith(".mp3"):
            audio_bytes = open(file_path, "rb").read()
            transcript = transcribe_audio(audio_bytes)
            raw_input_text = transcript
            input_type = "audio"

        else:
            return JSONResponse(status_code=400, content={"error": "Unsupported file type"})

    if not raw_input_text:
        return JSONResponse(status_code=400, content={"error": "No input provided"})

    # 3️⃣ Generate PRD JSON
    prd_json = generate_prd(raw_input_text)





   ###

    from datetime import datetime, timedelta

    # 1️⃣ Set correct PRD creation date
    current_date = datetime.utcnow().strftime("%d/%m/%Y")
    prd_json["date"] = current_date

    # 2️⃣ Fix ETA logic
    eta_value = prd_json.get("eta")

    if not eta_value or "90" in str(eta_value).lower():
        eta_date = datetime.utcnow() + timedelta(days=90)
        prd_json["eta"] = eta_date.strftime("%d/%m/%Y")

   ###








    # 4️⃣ Prepare Data for MongoDB
    db_data = {
        "project_name": prd_json.get("project_name"),
        "date": prd_json.get("date"),
        "requirements": prd_json.get("requirements"),
        "team_tasks": prd_json.get("team_tasks"),
        "eta": prd_json.get("eta"),
        "project_overview": prd_json.get("project_overview"),
        "raw_input": raw_input_text,
        "transcript": transcript,
        "input_type": input_type
    }

    inserted_id = save_prd(db_data)

    return {
        "message": "PRD Generated Successfully",
    #     "mongo_id": inserted_id,
        "data": prd_json
    }
